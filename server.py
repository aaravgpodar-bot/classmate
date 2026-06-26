import json
import io
import hmac
import os
import re
import sqlite3
import urllib.error
import urllib.parse
import urllib.request
from contextlib import contextmanager
from pathlib import Path

from flask import Flask, jsonify, request, send_file, send_from_directory, session

from ai_helpers import ai_extract_timetable, ai_game_questions, ai_paraphrase, ai_presentation_plan


BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = Path(os.environ.get("CLASSMATE_DATA_DIR", BASE_DIR / "data"))
DB_PATH = DATA_DIR / "classmate.sqlite3"
PUBLIC_FILES = {
    "app.js",
    "index.html",
    "manifest.json",
    "service-worker.js",
    "styles.css",
}

app = Flask(__name__, static_folder=str(BASE_DIR), static_url_path="")
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "classmate-dev-secret-change-me")
app.config.update(
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE="Lax",
    SESSION_COOKIE_SECURE=os.environ.get("FLASK_COOKIE_SECURE", "0") == "1",
)
DATA_DIR.mkdir(parents=True, exist_ok=True)


@app.get("/")
def index():
    return send_from_directory(BASE_DIR, "index.html")


@app.get("/health")
def health():
    return jsonify({"status": "ok", "app": "ClassMate"})


@app.get("/api/config")
def config():
    return jsonify({"googleClientId": os.environ.get("GOOGLE_CLIENT_ID", "")})


@app.post("/api/google-login")
def google_login():
    body = request.get_json(silent=True) or {}
    credential = str(body.get("credential", "")).strip()
    client_id = os.environ.get("GOOGLE_CLIENT_ID", "")
    if not client_id:
        return jsonify({"error": "GOOGLE_CLIENT_ID is not configured on the server."}), 503
    if not credential:
        return jsonify({"error": "Google credential is required."}), 400
    try:
        user = verify_google_credential(credential, client_id)
    except RuntimeError as error:
        return jsonify({"error": str(error)}), 401
    session["workspace_id"] = f"google:{user['email'].lower()}"
    session["email"] = user["email"].lower()
    session["role"] = str(body.get("role", "student")).strip()[:20] or "student"
    return jsonify({"user": user, "workspaceId": session["workspace_id"]})


@app.post("/api/logout")
def logout():
    session.clear()
    return jsonify({"signedOut": True})


@app.get("/api/workspace/<workspace_id>")
def get_workspace(workspace_id):
    workspace_id = clean_workspace_id(workspace_id)
    if not workspace_id:
        return jsonify({"error": "Workspace id is required."}), 400
    auth_error = authorize_workspace(workspace_id, allow_create=True)
    if auth_error:
        return auth_error
    saved = load_workspace(workspace_id)
    return jsonify({"workspaceId": workspace_id, "state": saved})


@app.post("/api/workspace/<workspace_id>")
def save_workspace(workspace_id):
    workspace_id = clean_workspace_id(workspace_id)
    if not workspace_id:
        return jsonify({"error": "Workspace id is required."}), 400
    body = request.get_json(silent=True) or {}
    state = body.get("state", {})
    if not isinstance(state, dict):
        return jsonify({"error": "Workspace state must be an object."}), 400
    auth_error = authorize_workspace(workspace_id, allow_create=True)
    if auth_error:
        return auth_error
    save_workspace_state(workspace_id, state)
    return jsonify({"workspaceId": workspace_id, "saved": True})


@app.delete("/api/workspace/<workspace_id>")
def delete_workspace(workspace_id):
    workspace_id = clean_workspace_id(workspace_id)
    if not workspace_id:
        return jsonify({"error": "Workspace id is required."}), 400
    auth_error = authorize_workspace(workspace_id, allow_create=False)
    if auth_error:
        return auth_error
    with db() as connection:
        connection.execute("DELETE FROM workspaces WHERE workspace_id = ?", (workspace_id,))
    return jsonify({"workspaceId": workspace_id, "deleted": True})


@app.post("/api/generate-game")
def generate_game():
    body = request.get_json(silent=True) or {}
    subject = str(body.get("subject", "")).strip()[:80] or "General Knowledge"
    try:
        questions = ai_game_questions(subject)
    except RuntimeError as error:
        return jsonify({"error": str(error)}), 503
    return jsonify({"subject": subject, "source": "openai", "questions": questions})


@app.post("/api/generate-presentation")
def generate_presentation():
    body = request.get_json(silent=True) or {}
    topic = str(body.get("topic", "")).strip()[:120] or "Untitled presentation"
    grade = str(body.get("grade", "")).strip()[:60] or "student"
    style = str(body.get("style", "")).strip()[:80] or "creative but school-ready"
    try:
        plan = ai_presentation_plan(topic, grade, style)
    except RuntimeError as error:
        return jsonify({"error": str(error)}), 503
    return jsonify({"topic": topic, "source": "openai", "plan": plan})


@app.post("/api/paraphrase")
def paraphrase():
    body = request.get_json(silent=True) or {}
    text = str(body.get("text", "")).strip()[:5000]
    tone = str(body.get("tone", "")).strip()[:80] or "clear student submission"
    if not text:
        return jsonify({"error": "Text is required for paraphrasing."}), 400
    try:
        result = ai_paraphrase(text, tone)
    except RuntimeError as error:
        return jsonify({"error": str(error)}), 503
    return jsonify({"source": "openai", "result": result})


@app.post("/api/extract-timetable")
def extract_timetable():
    body = request.get_json(silent=True) or {}
    image = str(body.get("image", "")).strip()
    if not image.startswith("data:image/"):
        return jsonify({"error": "A timetable image is required."}), 400
    try:
        result = ai_extract_timetable(image[:8_000_000])
    except RuntimeError as error:
        return jsonify({"error": str(error)}), 503
    return jsonify({"source": "openai", "timetable": result})


@app.post("/api/export-assignment")
def export_assignment():
    body = request.get_json(silent=True) or {}
    assignment = body.get("assignment", {}) if isinstance(body.get("assignment"), dict) else {}
    title = str(assignment.get("title", "ClassMate project")).strip()[:120] or "ClassMate project"
    kind = str(assignment.get("kind", "PowerPoint PPTX")).strip()[:80] or "PowerPoint PPTX"
    sections = assignment.get("sections", []) if isinstance(assignment.get("sections"), list) else []
    file_base = safe_filename(title)
    try:
        if re.search(r"word|docx|notion|pdf", kind, flags=re.I):
            data = build_docx_export(title, kind, assignment, sections)
            return send_file(
                data,
                as_attachment=True,
                download_name=f"{file_base}.docx",
                mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        data = build_pptx_export(title, kind, assignment, sections)
        return send_file(
            data,
            as_attachment=True,
            download_name=f"{file_base}.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except ImportError:
        return jsonify({"error": "Export libraries are not installed. Run pip install -r requirements.txt."}), 503


@app.get("/<path:filename>")
def static_files(filename):
    if filename not in PUBLIC_FILES:
        return jsonify({"error": "not found"}), 404
    return send_from_directory(BASE_DIR, filename)


def init_db():
    with sqlite3.connect(DB_PATH) as connection:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS workspaces (
                workspace_id TEXT PRIMARY KEY,
                owner_kind TEXT NOT NULL DEFAULT 'guest',
                owner_id TEXT NOT NULL DEFAULT '',
                workspace_secret TEXT NOT NULL DEFAULT '',
                state_json TEXT NOT NULL,
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        columns = {row[1] for row in connection.execute("PRAGMA table_info(workspaces)").fetchall()}
        if "owner_kind" not in columns:
            connection.execute("ALTER TABLE workspaces ADD COLUMN owner_kind TEXT NOT NULL DEFAULT 'guest'")
        if "owner_id" not in columns:
            connection.execute("ALTER TABLE workspaces ADD COLUMN owner_id TEXT NOT NULL DEFAULT ''")
        if "workspace_secret" not in columns:
            connection.execute("ALTER TABLE workspaces ADD COLUMN workspace_secret TEXT NOT NULL DEFAULT ''")


@contextmanager
def db():
    connection = sqlite3.connect(DB_PATH)
    try:
        yield connection
        connection.commit()
    except Exception:
        connection.rollback()
        raise
    finally:
        connection.close()


def workspace_owner(workspace_id):
    if workspace_id.startswith("google:"):
        return "google", workspace_id.split(":", 1)[1].lower()
    return "guest", workspace_id


def request_workspace_secret():
    return str(request.headers.get("X-ClassMate-Workspace-Secret", "")).strip()


def authorize_workspace(workspace_id, allow_create):
    owner_kind, owner_id = workspace_owner(workspace_id)
    if owner_kind == "google":
        if session.get("workspace_id") != workspace_id:
            return jsonify({"error": "Sign in with the matching Google account to access this workspace."}), 403
        return None

    secret = request_workspace_secret()
    if not secret:
        return jsonify({"error": "Workspace secret is required for guest sync."}), 403
    with db() as connection:
        row = connection.execute(
            "SELECT workspace_secret FROM workspaces WHERE workspace_id = ?",
            (workspace_id,),
        ).fetchone()
    if not row:
        return None if allow_create else (jsonify({"error": "Workspace not found."}), 404)
    if not hmac.compare_digest(row[0] or "", secret):
        return jsonify({"error": "This device is not allowed to access that guest workspace."}), 403
    return None


init_db()


def clean_workspace_id(value):
    return re.sub(r"[^a-zA-Z0-9_.:@-]", "", str(value or "").strip())[:160]


def load_workspace(workspace_id):
    with db() as connection:
        row = connection.execute(
            "SELECT state_json FROM workspaces WHERE workspace_id = ?",
            (workspace_id,),
        ).fetchone()
    if not row:
        return None
    try:
        return json.loads(row[0])
    except json.JSONDecodeError:
        return None


def save_workspace_state(workspace_id, state):
    state_json = json.dumps(state, separators=(",", ":"), ensure_ascii=False)
    owner_kind, owner_id = workspace_owner(workspace_id)
    workspace_secret = request_workspace_secret() if owner_kind == "guest" else ""
    with db() as connection:
        connection.execute(
        """
            INSERT INTO workspaces (workspace_id, owner_kind, owner_id, workspace_secret, state_json, updated_at)
            VALUES (?, ?, ?, ?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(workspace_id) DO UPDATE SET
              owner_kind = excluded.owner_kind,
              owner_id = excluded.owner_id,
              workspace_secret = CASE
                WHEN workspaces.workspace_secret = '' THEN excluded.workspace_secret
                ELSE workspaces.workspace_secret
              END,
              state_json = excluded.state_json,
              updated_at = CURRENT_TIMESTAMP
            """,
            (workspace_id, owner_kind, owner_id, workspace_secret, state_json),
        )


def verify_google_credential(credential, client_id):
    url = "https://oauth2.googleapis.com/tokeninfo?id_token=" + urllib.parse.quote(credential)
    try:
      with urllib.request.urlopen(url, timeout=10) as response:
          payload = json.loads(response.read().decode("utf-8"))
    except (urllib.error.URLError, json.JSONDecodeError) as error:
      raise RuntimeError("Could not verify Google sign-in.") from error
    if payload.get("aud") != client_id:
      raise RuntimeError("Google sign-in was not issued for this ClassMate app.")
    if payload.get("email_verified") not in (True, "true", "True"):
      raise RuntimeError("Google email is not verified.")
    return {
      "name": payload.get("name") or payload.get("email", "").split("@")[0] or "Student",
      "email": payload.get("email", ""),
      "picture": payload.get("picture", ""),
    }


def safe_filename(value):
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()
    return cleaned[:60] or "classmate-project"


def section_text(section):
    name = str(section.get("name", "Section")).strip()
    status = str(section.get("status", "To do")).strip()
    owner = str(section.get("owner", "You")).strip()
    comments = str(section.get("comments", 0))
    suggestions = str(section.get("suggestions", 0))
    return name, f"Status: {status}\nOwner: {owner}\nComments: {comments}\nSuggested edits: {suggestions}"


def build_pptx_export(title, kind, assignment, sections):
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = f"ClassMate template for {kind}"
    for section in sections[:12] or [{"name": "Plan", "status": "To do", "owner": "You"}]:
        name, body = section_text(section)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = name
        textbox = slide.placeholders[1]
        textbox.text = body
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
    if assignment.get("theme"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Visual Direction"
        slide.placeholders[1].text = str(assignment.get("theme", "Clean student style"))
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def build_docx_export(title, kind, assignment, sections):
    from docx import Document

    document = Document()
    document.add_heading(title, 0)
    document.add_paragraph(f"ClassMate template for {kind}")
    if assignment.get("theme"):
        document.add_paragraph(f"Look: {assignment.get('theme')}")
    document.add_paragraph(f"Lead: {assignment.get('lead', 'You')}")
    document.add_paragraph(f"Feedback: {assignment.get('feedback', 'Balanced')}")
    for index, section in enumerate(sections[:30] or [{"name": "Plan", "status": "To do", "owner": "You"}], start=1):
        name, body = section_text(section)
        document.add_heading(f"{index}. {name}", level=1)
        for line in body.splitlines():
            document.add_paragraph(line)
    output = io.BytesIO()
    document.save(output)
    output.seek(0)
    return output



if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5177, debug=True)

